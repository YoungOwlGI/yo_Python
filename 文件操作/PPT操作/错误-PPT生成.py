# 一个做PPT的错误程序

from pptx import Presentation
from pptx.util import Inches

# 创建一个新的PPT对象
presentation = Presentation()

# 添加一个版画的标题幻灯片
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "版画艺术欣赏"
subtitle.text = "介绍版画的制作过程和艺术价值"

# 添加一个版画的介绍幻灯片
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)

left = Inches(1)
top = Inches(1)
width = Inches(6)
height = Inches(3)

# 添加图片
slide.shapes.add_picture(r'image.jpg', left, top, width, height)

# 添加文本框
txtbox = slide.shapes.add_textbox(left, top, width, height)
p = txtbox.paragraphs[0]
p.text = "这里是版画的介绍文本，可以添加版画的制作材料、历史、流派等信息。"

# 保存PPT文件
presentation.save('版画艺术欣赏.pptx')